
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, io, os

# ── Colors ─────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
PINK    = RGBColor(0xE9, 0x5B, 0x8C)   # rose-pink (accent)
PINK_LT = RGBColor(0xFF, 0xF0, 0xF5)   # light pink bg
GRAY    = RGBColor(0x55, 0x55, 0x66)   # body text
BLUE    = RGBColor(0x2C, 0x3E, 0x7A)   # heading dark blue
GREEN   = RGBColor(0x27, 0xAE, 0x60)

# ── Slide dimensions ───────────────────────────────────
W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_layout(prs):
    return prs.slide_layouts[6]   # blank

def add_rect(slide, x, y, w, h, fill=None, border=None, radius=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb

def bg(slide, color=PINK_LT):
    """Fill slide background."""
    bg_shape = add_rect(slide, 0, 0, W, H, fill=color)
    # send to back
    sp_tree = slide.shapes._spTree
    sp_tree.remove(bg_shape._element)
    sp_tree.insert(2, bg_shape._element)

def add_picture_from_path(slide, img_path, x, y, w, h):
    slide.shapes.add_picture(img_path, x, y, w, h)

def top_bar(slide, title_text, subtitle=""):
    """White top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.4), fill=WHITE)
    # Pink left accent strip
    add_rect(slide, 0, 0, Inches(0.18), Inches(1.4), fill=PINK)
    add_text(slide, title_text, Inches(0.35), Inches(0.18), Inches(10), Inches(0.8),
             size=28, bold=True, color=BLUE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.90), Inches(10), Inches(0.4),
                 size=13, color=GRAY)

MEDIA = r"d:\Best cook\presentation\extracted\ppt\media"

# ════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl, WHITE)
# Pink left half
add_rect(sl, 0, 0, Inches(5.8), H, fill=PINK)

add_text(sl, "BEST CAKE", Inches(0.5), Inches(1.8), Inches(5), Inches(1.4),
         size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "AI yordamida tort buyurtma qilish platformasi",
         Inches(0.5), Inches(3.4), Inches(5), Inches(0.7),
         size=17, color=WHITE)

# 3 feature pills
features = [("🎨", "Dizayn yarating"),
            ("📦", "3D ko'ring"),
            ("🛒", "Tez buyurtma")]
for i, (ic, lb) in enumerate(features):
    fy = Inches(4.3) + i * Inches(0.72)
    add_rect(sl, Inches(0.5), fy, Inches(4.8), Inches(0.56),
             fill=RGBColor(0xFF,0xFF,0xFF), border=None)
    add_text(sl, f"{ic}  {lb}", Inches(0.65), fy + Inches(0.06), Inches(4.5), Inches(0.45),
             size=14, bold=True, color=PINK)

# Right side cake image
add_picture_from_path(sl, os.path.join(MEDIA,"image1.png"),
                       Inches(5.7), Inches(0.2), Inches(7.3), Inches(7.1))

# ════════════════════════════════════════════════════════
# SLIDE 2 — Muammo (Problem)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "MUAMMO", "Nega tort buyurtma qilish hanuz qiyin?")

problems = [
    ("😕", "Mijoz va qandolatchi bir-birini tushunmaydi"),
    ("⚠️", "Allergiya chaqiruvchi ingredientlar nazorat qilinmaydi"),
    ("⏰", "Buyurtma berish ko'p vaqt oladi"),
]
for i, (ic, txt) in enumerate(problems):
    py = Inches(1.8) + i * Inches(1.4)
    add_rect(sl, Inches(0.5), py, Inches(6.5), Inches(1.1),
             fill=WHITE, border=RGBColor(0xE0,0xE0,0xE0))
    add_text(sl, ic, Inches(0.7), py + Inches(0.2), Inches(0.5), Inches(0.7),
             size=22, color=PINK)
    add_text(sl, txt, Inches(1.4), py + Inches(0.2), Inches(5.3), Inches(0.7),
             size=15, color=BLACK)

# Result box
add_rect(sl, Inches(0.5), Inches(5.8), Inches(6.5), Inches(1.3),
         fill=PINK, border=None)
add_text(sl, "Natija: Norozi mijozlar, isrof vaqt, qo'shimcha xarajatlar",
         Inches(0.7), Inches(5.95), Inches(6.1), Inches(1.0),
         size=14, bold=True, color=WHITE)

# Right image
add_picture_from_path(sl, os.path.join(MEDIA,"image3.png"),
                       Inches(7.2), Inches(1.5), Inches(5.7), Inches(5.5))

# ════════════════════════════════════════════════════════
# SLIDE 3 — Yechim (Solution)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "YECHIM", "AI + 3D + Tez buyurtma tizimi")

solutions = [
    ("🤖", "AI tavsiyalari", "Dizayn, ta'm va ingredientlar bo'yicha aqlli maslahat"),
    ("📦", "3D ko'rish",     "Tortni buyurtmadan oldin 3D formatda ko'rib chiqing"),
    ("🛒", "Tez buyurtma",   "Bir necha soniyada buyurtma bering"),
]
for i, (ic, title, desc) in enumerate(solutions):
    sy = Inches(1.8) + i * Inches(1.55)
    add_rect(sl, Inches(0.5), sy, Inches(6.5), Inches(1.3),
             fill=WHITE, border=RGBColor(0xE8,0xE8,0xE8))
    add_text(sl, ic, Inches(0.7), sy + Inches(0.15), Inches(0.6), Inches(0.8),
             size=24, color=PINK)
    add_text(sl, title, Inches(1.5), sy + Inches(0.1), Inches(5), Inches(0.45),
             size=15, bold=True, color=BLUE)
    add_text(sl, desc, Inches(1.5), sy + Inches(0.55), Inches(5), Inches(0.6),
             size=12, color=GRAY)

# Result
add_rect(sl, Inches(0.5), Inches(6.4), Inches(6.5), Inches(0.75),
         fill=BLUE, border=None)
add_text(sl, "✅  Natija: Aniq va qulay tort buyurtmasi",
         Inches(0.7), Inches(6.5), Inches(6.0), Inches(0.6),
         size=14, bold=True, color=WHITE)

add_picture_from_path(sl, os.path.join(MEDIA,"image4.png"),
                       Inches(7.2), Inches(1.5), Inches(5.7), Inches(5.5))

# ════════════════════════════════════════════════════════
# SLIDE 4 — Qanday Ishlaydi (How it works)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "QANDAY ISHLAYDI?", "Orzuingizdagi tortni 5 qadamda yarating")

steps = [
    ("1", "Tort turini\ntanlang"),
    ("2", "Ta'mlarni\ntanlang"),
    ("3", "Dizayn\nyarating"),
    ("4", "AI maslahat\noling"),
    ("5", "3D ko'ring va\nbuyurtma bering"),
]
sw = Inches(2.3)
for i, (num, label) in enumerate(steps):
    sx = Inches(0.35) + i * Inches(2.55)
    sy = Inches(1.8)
    # circle number
    add_rect(sl, sx, sy, Inches(0.55), Inches(0.55), fill=PINK)
    add_text(sl, num, sx, sy, Inches(0.55), Inches(0.55),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # box
    add_rect(sl, sx, sy + Inches(0.6), sw, Inches(2.8),
             fill=WHITE, border=RGBColor(0xDD,0xDD,0xDD))
    add_text(sl, label, sx + Inches(0.1), sy + Inches(0.75), sw - Inches(0.2), Inches(0.9),
             size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    # Arrow (not last)
    if i < 4:
        add_text(sl, "→", sx + sw + Inches(0.05), sy + Inches(1.6),
                 Inches(0.4), Inches(0.5), size=20, bold=True, color=PINK,
                 align=PP_ALIGN.CENTER)

# Bottom summary
add_rect(sl, Inches(0.35), Inches(5.9), Inches(12.5), Inches(1.2),
         fill=PINK, border=None)
add_text(sl, "Tasavvur qiling → Yarating → Bir necha soniyada buyurtma bering 🚀",
         Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.9),
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 5 — Raqobatchilar (Competitors)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "RAQOBATCHILAR TAQQOSLASH", "Boshgalar qisman yechim – BEST CAKE to'liq!")

headers = ["Xususiyat", "BEST CAKE", "Midjourney/DALL-E", "An'anaviy qandolatchilar"]
col_w   = [Inches(3.0), Inches(2.7), Inches(3.0), Inches(3.5)]
col_x   = [Inches(0.35), Inches(3.4), Inches(6.15), Inches(9.2)]
rows = [
    ("Dizayn ko'rish",        "✅ 3D real vaqt",         "✅ Rasm",         "❌ Yo'q"),
    ("AI tavsiyalari",        "✅ Dizayn+ta'm+allergiya","⚠️ Faqat rasm",   "❌ Yo'q"),
    ("Buyurtma tizimi",       "✅ To'liq avtomatik",     "❌ Yo'q",         "⚠️ Qo'lda"),
    ("Allergiya nazorati",    "✅ Mavjud",               "❌ Yo'q",         "❌ Yo'q"),
    ("Shaxsiylashtirish",     "✅ End-to-end",           "⚠️ Vizual",       "⚠️ Cheklangan"),
]

col_colors = [RGBColor(0xF5,0xF5,0xF5), PINK, WHITE, WHITE]
for ci, (hdr, cx, cw, cc) in enumerate(zip(headers, col_x, col_w, col_colors)):
    hy = Inches(1.7)
    hh = Inches(0.5)
    add_rect(sl, cx, hy, cw - Inches(0.05), hh,
             fill=BLUE if ci==1 else RGBColor(0x3A,0x3A,0x5A),
             border=None)
    add_text(sl, hdr, cx + Inches(0.05), hy, cw - Inches(0.1), hh,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        ry = Inches(2.25) + ri * Inches(0.88)
        rc = PINK_LT if ci == 1 else (WHITE if ri%2==0 else RGBColor(0xFB,0xFB,0xFB))
        add_rect(sl, cx, ry, cw - Inches(0.05), Inches(0.82),
                 fill=rc, border=RGBColor(0xE0,0xE0,0xE0))
        add_text(sl, row[ci], cx + Inches(0.08), ry + Inches(0.1),
                 cw - Inches(0.15), Inches(0.65),
                 size=12, color=BLUE if ci==1 else BLACK, bold=(ci==1))

# ════════════════════════════════════════════════════════
# SLIDE 6 — Bozor (Market)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "BOZOR IMKONIYATI", "Best Cake potensiali")

markets = [
    ("TAM", "Global bozor", "$150+ mlrd/yil", BLUE),
    ("SAM", "O'zbekiston",  "$500M+",          PINK),
    ("SOM", "Bizning ulush","$50–75M/yil",      GREEN),
]
for i, (tag, label, value, color) in enumerate(markets):
    mx = Inches(0.5) + i * Inches(4.15)
    my = Inches(1.8)
    add_rect(sl, mx, my, Inches(3.8), Inches(3.2),
             fill=WHITE, border=RGBColor(0xDD,0xDD,0xDD))
    add_rect(sl, mx, my, Inches(3.8), Inches(0.55), fill=color)
    add_text(sl, tag, mx + Inches(0.1), my + Inches(0.06), Inches(1.0), Inches(0.45),
             size=18, bold=True, color=WHITE)
    add_text(sl, label, mx + Inches(0.1), my + Inches(0.75), Inches(3.5), Inches(0.5),
             size=14, color=GRAY)
    add_text(sl, value, mx + Inches(0.1), my + Inches(1.4), Inches(3.5), Inches(0.8),
             size=26, bold=True, color=color)

# Bottom 3 problems
add_text(sl, "Hozirgi asosiy muammolar:", Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE)
pain = ["❌ AI/3D yechim yo'q", "❌ Allergiya nazorati yo'q", "❌ Buyurtma murakkab"]
for i, p in enumerate(pain):
    px = Inches(0.5) + i * Inches(4.15)
    add_rect(sl, px, Inches(5.9), Inches(3.9), Inches(0.9),
             fill=RGBColor(0xFF,0xEE,0xEE), border=RGBColor(0xFF,0xCC,0xCC))
    add_text(sl, p, px + Inches(0.1), Inches(5.95), Inches(3.7), Inches(0.7),
             size=13, color=BLACK)

# ════════════════════════════════════════════════════════
# SLIDE 7 — Daromad Modeli (Business Model)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "DAROMAD MODELI", "Best Cake qanday pul ishlaydi?")

models = [
    ("1", "Buyurtma komissiyasi",   "Tortlardan 10%\nFast food dan 5%",        PINK),
    ("2", "Premium hamkorlik",       "Qandolatchi uchun 'TOP LIST'\nKo'proq ko'rinish = ko'proq buyurtma", BLUE),
    ("3", "AI Analytics xizmatlari","Sotuv, talab va trend tahlillari\nQandolatchilar uchun biznes maslahat", GREEN),
    ("4", "Reklama va Boost tizimi", "Mahsulotni oldinga chiqarish\nMaqsadli reklama joylash",  RGBColor(0xF3,0x9C,0x12)),
]
for i, (num, title, desc, color) in enumerate(models):
    row = i // 2
    col = i % 2
    mx = Inches(0.5) + col * Inches(6.25)
    my = Inches(1.8) + row * Inches(2.4)
    add_rect(sl, mx, my, Inches(6.0), Inches(2.1),
             fill=WHITE, border=RGBColor(0xE0,0xE0,0xE0))
    add_rect(sl, mx, my, Inches(0.55), Inches(2.1), fill=color)
    add_text(sl, num, mx + Inches(0.05), my + Inches(0.65), Inches(0.45), Inches(0.6),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, mx + Inches(0.7), my + Inches(0.15), Inches(5.1), Inches(0.5),
             size=14, bold=True, color=BLACK)
    add_text(sl, desc, mx + Inches(0.7), my + Inches(0.7), Inches(5.1), Inches(1.2),
             size=12, color=GRAY)

# Bottom
add_rect(sl, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), fill=BLUE)
add_text(sl, "\"Biz faqat buyurtma qilmaymiz — biznesni o'stiramiz!\"",
         Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.6),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 8 — Milestones & KPI
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "MILESTONES & KPI", "Best Cake rivojlanish yo'li")

phases = [
    ("🚀 MVP  (0–1 oy)",   ["100–300 foydalanuvchi","10–20 hamkor","300–700 buyurtma","$300–$1,000 daromad"], PINK),
    ("📈 6 OY",            ["5,000–10,000 foydalanuvchi","100+ hamkor","15,000+ buyurtma","$10K–$25K daromad"], BLUE),
    ("🌍 3 YIL",           ["200,000+ foydalanuvchi","1,000+ hamkor","1,000,000+ buyurtma","$1M+ yillik daromad"], GREEN),
]
for i, (phase, kpis, color) in enumerate(phases):
    px = Inches(0.4) + i * Inches(4.25)
    py = Inches(1.7)
    add_rect(sl, px, py, Inches(4.0), Inches(0.6), fill=color)
    add_text(sl, phase, px + Inches(0.1), py + Inches(0.05), Inches(3.8), Inches(0.5),
             size=14, bold=True, color=WHITE)
    add_rect(sl, px, py + Inches(0.6), Inches(4.0), Inches(4.3),
             fill=WHITE, border=RGBColor(0xDD,0xDD,0xDD))
    for ki, kpi in enumerate(kpis):
        ky = py + Inches(0.8) + ki * Inches(0.88)
        add_rect(sl, px + Inches(0.12), ky, Inches(3.75), Inches(0.72),
                 fill=PINK_LT if i==0 else (RGBColor(0xEE,0xF4,0xFF) if i==1 else RGBColor(0xEE,0xFF,0xF2)),
                 border=None)
        add_text(sl, "• " + kpi, px + Inches(0.22), ky + Inches(0.1), Inches(3.5), Inches(0.55),
                 size=12, color=BLACK)

# Bottom
add_rect(sl, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.8), fill=BLUE)
add_text(sl, "BEST CAKE — kichik MVPdan global AI marketplacegacha o'sadigan platforma 🚀",
         Inches(0.6), Inches(6.48), Inches(12.0), Inches(0.65),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 9 — Jamoa (Team)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
bg(sl)
top_bar(sl, "BIZNING JAMOA", "Best Cake ortidagi odamlar")

team = [
    (os.path.join(MEDIA,"image11.png"), "Umidjon Odiljonov",  "CEO / Loyiha rahbari"),
    (os.path.join(MEDIA,"image14.png"), "Sardorbek Albakiev", "CTO / Dasturchi"),
    (os.path.join(MEDIA,"image12.png"), "Abdulloh Ma'murov",  "Co-founder"),
    (os.path.join(MEDIA,"image13.png"), "Bogdan Stadnyuk",    "Biznes analitik"),
]
card_colors = [PINK, BLUE, GREEN, RGBColor(0xF3,0x9C,0x12)]
for i, (img, name, role) in enumerate(team):
    tx = Inches(0.4) + i * Inches(3.2)
    ty = Inches(1.7)
    tw = Inches(3.0)
    # Card
    add_rect(sl, tx, ty, tw, Inches(5.4), fill=WHITE,
             border=RGBColor(0xE0,0xE0,0xE0))
    # Top color accent
    add_rect(sl, tx, ty, tw, Inches(0.25), fill=card_colors[i])
    # Photo
    try:
        add_picture_from_path(sl, img, tx + Inches(0.5), ty + Inches(0.4), Inches(2.0), Inches(2.7))
    except Exception as e:
        print(f"Photo error {img}: {e}")
    # Name
    add_text(sl, name, tx + Inches(0.1), ty + Inches(3.3), tw - Inches(0.2), Inches(0.55),
             size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    # Role badge
    add_rect(sl, tx + Inches(0.3), ty + Inches(3.95), tw - Inches(0.6), Inches(0.5),
             fill=card_colors[i])
    add_text(sl, role, tx + Inches(0.3), ty + Inches(3.97), tw - Inches(0.6), Inches(0.45),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out_path = r"d:\Best cook\presentation\BestCake_Simple.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
